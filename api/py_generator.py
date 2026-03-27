import sys
import json
import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as WordPt, Inches as WordInches, RGBColor as WordRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import pandas as pd
from io import BytesIO
import tempfile
import matplotlib.pyplot as plt
import seaborn as sns
from duckduckgo_search import DDGS
from bs4 import BeautifulSoup
import base64

def web_search(query, max_results=5):
    try:
        results = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                results.append({
                    "title": r.get("title", ""),
                    "url": r.get("href", ""),
                    "body": r.get("body", "")
                })
        return results
    except Exception as e:
        return f"Search Error: {str(e)}"

def generate_excel(data, filename):
    try:
        cols = data.get("columns", [])
        rows = data.get("rows", [])
        df = pd.DataFrame(rows, columns=cols if cols else None)
        output_path = os.path.join("/tmp", filename)
        writer = pd.ExcelWriter(output_path, engine='xlsxwriter')
        df.to_excel(writer, index=False, sheet_name='Veri')
        workbook = writer.book
        worksheet = writer.sheets['Veri']
        header_format = workbook.add_format({
            'bold': True, 'text_wrap': True, 'valign': 'top',
            'fg_color': '#2DD4BF', 'font_color': '#FFFFFF', 'border': 1
        })
        for col_num, value in enumerate(df.columns.values):
            worksheet.write(0, col_num, value, header_format)
        writer.close()
        
        # Preview ve kod snippet oluştur
        code_snippet = generate_code_snippet(data, filename)
        preview_image = generate_excel_preview_image(data, filename)
        
        return {
            "path": output_path,
            "code_snippet": code_snippet,
            "preview_image": preview_image,
            "rows": len(df),
            "columns": len(df.columns)
        }
    except Exception as e:
        return f"Error: {str(e)}"

def generate_code_snippet(data, filename):
    """Excel oluşturmak için kullanılan Python kodunu döndür"""
    cols = data.get("columns", [])
    rows = data.get("rows", [])
    
    code = f"""import pandas as pd

# Veri hazırlama
columns = {repr(cols)}
rows = {repr(rows[:5] if len(rows) > 5 else rows)}...

# DataFrame oluşturma
df = pd.DataFrame(rows, columns=columns)

# Excel dosyasına kaydetme
output_path = "/tmp/{filename}"
writer = pd.ExcelWriter(output_path, engine='xlsxwriter')
df.to_excel(writer, index=False, sheet_name='Veri')

# Formatlama
workbook = writer.book
worksheet = writer.sheets['Veri']
header_format = workbook.add_format({
    'bold': True, 'fg_color': '#2DD4BF',
    'font_color': '#FFFFFF', 'border': 1
})
writer.close()
print(f"✅ {{len(df)}} satır, {{len(columns)}} sütun oluşturuldu")
"""
    return code

def generate_excel_preview_image(data, filename):
    """Excel dosyasının görsel önizlemesini oluştur"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        from matplotlib import font_manager as fm
        import matplotlib.pyplot as plt
        from matplotlib.patches import Rectangle
        
        cols = data.get("columns", [])
        rows = data.get("rows", [])
        
        if not cols and not rows:
            return None
            
        # DataFrame oluştur
        df = pd.DataFrame(rows, columns=cols if cols else None)
        
        # Görsel boyutları
        n_rows = min(len(df), 15)  # Maksimum 15 satır göster
        n_cols = len(df.columns)
        
        fig_height = max(4, n_rows * 0.4 + 1.5)
        fig_width = max(8, n_cols * 1.5)
        
        fig, ax = plt.subplots(figsize=(fig_width, fig_height))
        ax.axis('off')
        
        # Başlık
        fig.suptitle(f'📊 {filename}', fontsize=14, fontweight='bold', 
                     color='#2DD4BF', y=0.98)
        
        # Tablo oluştur
        cell_height = 0.08
        cell_width = 1.5
        start_y = 0.88
        
        # Header arka planı
        header_rect = Rectangle((0.05, start_y - cell_height), 
                                min(n_cols, 8) * cell_width + 0.1, 
                                cell_height, 
                                facecolor='#2DD4BF', edgecolor='none', 
                                transform=fig.transFigure)
        ax.add_patch(header_rect)
        
        # Sütun başlıkları
        for i, col in enumerate(df.columns[:8]):  # Maksimum 8 sütun
            x = 0.1 + i * cell_width
            ax.text(x, start_y - cell_height/2, str(col)[:20], 
                   fontsize=9, fontweight='bold', color='white',
                   transform=fig.transFigure, va='center')
        
        # Veri satırları
        for row_idx in range(n_rows):
            y = start_y - (row_idx + 1) * cell_height - 0.02
            bg_color = '#1E293B' if row_idx % 2 == 0 else '#0F172A'
            
            # Satır arka planı
            row_rect = Rectangle((0.05, y - cell_height/2), 
                                min(n_cols, 8) * cell_width + 0.1, 
                                cell_height, 
                                facecolor=bg_color, edgecolor='#334155', 
                                linewidth=0.5,
                                transform=fig.transFigure)
            ax.add_patch(row_rect)
            
            # Hücre verileri
            for col_idx in range(min(n_cols, 8)):
                x = 0.1 + col_idx * cell_width
                val = str(df.iloc[row_idx, col_idx])[:25]
                ax.text(x, y, val, fontsize=8, color='#E2E8F0',
                       transform=fig.transFigure, va='center')
        
        # Alt bilgi
        footer_text = f"Toplam {len(df)} satır | {len(df.columns)} sütun | 💾 İndirmeye hazır"
        ax.text(0.5, 0.05, footer_text, fontsize=9, color='#94A3B8',
               transform=fig.transFigure, ha='center')
        
        # Kenar boşlukları
        fig.patch.set_facecolor('#0F172A')
        ax.set_facecolor('#0F172A')
        
        # Görseli kaydet
        preview_path = os.path.join("/tmp", f"preview_{filename}.png")
        plt.savefig(preview_path, dpi=150, bbox_inches='tight', 
                   facecolor='#0F172A', edgecolor='none')
        plt.close()
        
        # Base64'e çevir
        with open(preview_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode('utf-8')
        
        os.remove(preview_path)
        return img_data
        
    except Exception as e:
        print(f"Preview image error: {e}", file=sys.stderr)
        return None

def generate_word(content, filename):
    try:
        doc = Document()
        doc.add_heading('LifeCoach AI - Rapor', 0)
        if isinstance(content, str):
            for p in content.split('\n\n'):
                doc.add_paragraph(p)
        elif isinstance(content, list):
            for section in content:
                if section.get("type") == "heading":
                    doc.add_heading(section.get("text", ""), level=section.get("level", 1))
                elif section.get("type") == "image":
                    # For EQ charts
                    img_path = section.get("path")
                    if img_path and os.path.exists(img_path):
                        doc.add_picture(img_path, width=Inches(5))
                else:
                    doc.add_paragraph(section.get("text", ""))
        output_path = os.path.join("/tmp", filename)
        doc.save(output_path)
        return output_path
    except Exception as e:
        return f"Error: {str(e)}"

def _existing_file(paths):
    for p in paths:
        if p and os.path.exists(p):
            return p
    return None

def generate_certificate(data, filename):
    try:
        user_name = (data.get("userName") or "Değerli Kullanıcı").strip()
        goal_title = (data.get("goalTitle") or "Belirlenen hedef").strip()
        issue_date = (data.get("date") or "").strip()
        issuer_name = (data.get("issuerName") or "METEHAN HAYDAR ERBAŞ").strip()
        issuer_title = (data.get("issuerTitle") or "Founder, HAN AI").strip()

        if not issue_date:
            from datetime import datetime
            issue_date = datetime.now().strftime("%d.%m.%Y")

        # Signature & logo inputs (optional, fallback friendly)
        provided_signature = data.get("signaturePath")
        provided_logo = data.get("logoPath")
        signature_path = _existing_file([
            provided_signature,
            "/home/spectre05/.cursor/projects/home-spectre05-Masa-st-LifeCoach-Cloude/assets/WhatsApp_Image_2026-03-22_at_22.25.40-4e96b3ad-6e64-4581-8977-a98945d0266a.png",
            os.path.join(os.getcwd(), "public", "han-signature.png"),
            os.path.join(os.getcwd(), "public", "signature.png"),
        ])
        logo_path = _existing_file([
            provided_logo,
            os.path.join(os.getcwd(), "public", "lifecoach_logo.png"),
            os.path.join(os.getcwd(), "public", "lifecoach-logo-splash.png"),
        ])

        doc = Document()
        section = doc.sections[0]
        section.top_margin = WordInches(0.6)
        section.bottom_margin = WordInches(0.6)
        section.left_margin = WordInches(0.7)
        section.right_margin = WordInches(0.7)

        # Top branding area
        if logo_path:
            p_logo = doc.add_paragraph()
            p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_logo.add_run().add_picture(logo_path, width=WordInches(1.2))

        p_brand = doc.add_paragraph("HAN AI | LIFECOACH")
        p_brand.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_brand = p_brand.runs[0]
        r_brand.bold = True
        r_brand.font.size = WordPt(13)
        r_brand.font.color.rgb = WordRGBColor(15, 118, 110)

        p_title = doc.add_paragraph("CERTIFICATE OF ACHIEVEMENT")
        p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_title = p_title.runs[0]
        r_title.bold = True
        r_title.font.size = WordPt(28)
        r_title.font.color.rgb = WordRGBColor(17, 24, 39)

        p_sub = doc.add_paragraph("This certificate is proudly presented to")
        p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_sub.runs[0].font.size = WordPt(12)

        p_name = doc.add_paragraph(user_name.upper())
        p_name.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_name = p_name.runs[0]
        r_name.bold = True
        r_name.font.size = WordPt(24)
        r_name.font.color.rgb = WordRGBColor(5, 150, 105)

        p_goal = doc.add_paragraph("for successfully completing the goal below:")
        p_goal.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_goal.runs[0].font.size = WordPt(11)

        p_goal_title = doc.add_paragraph(f"\"{goal_title}\"")
        p_goal_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_goal = p_goal_title.runs[0]
        r_goal.bold = True
        r_goal.font.size = WordPt(18)
        r_goal.font.color.rgb = WordRGBColor(3, 105, 161)

        p_stamp = doc.add_paragraph("HAN AI VERIFIED")
        p_stamp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r_stamp = p_stamp.runs[0]
        r_stamp.bold = True
        r_stamp.font.size = WordPt(14)
        r_stamp.font.color.rgb = WordRGBColor(185, 28, 28)

        doc.add_paragraph("")

        table = doc.add_table(rows=1, cols=2)
        table.autofit = True
        left_cell = table.rows[0].cells[0]
        right_cell = table.rows[0].cells[1]

        p_date_label = left_cell.paragraphs[0]
        p_date_label.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run_date_label = p_date_label.add_run("Date\n")
        run_date_label.bold = True
        run_date_label.font.size = WordPt(10)
        run_date = p_date_label.add_run(issue_date)
        run_date.font.size = WordPt(12)

        p_sign = right_cell.paragraphs[0]
        p_sign.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p_sign.add_run("Authorized by\n").bold = True
        if signature_path:
            p_sign.add_run().add_picture(signature_path, width=WordInches(1.8))
            p_sign.add_run("\n")
        r_issuer = p_sign.add_run(issuer_name)
        r_issuer.bold = True
        r_issuer.font.size = WordPt(11)
        p_sign.add_run(f"\n{issuer_title}").font.size = WordPt(9)

        footer = doc.add_paragraph("Official digital certificate generated by HAN AI LifeCoach.")
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer.runs[0].italic = True
        footer.runs[0].font.size = WordPt(9)
        footer.runs[0].font.color.rgb = WordRGBColor(100, 116, 139)

        output_path = os.path.join("/tmp", filename)
        doc.save(output_path)
        return output_path
    except Exception as e:
        return f"Error: {str(e)}"

def generate_eq_chart(mood_history, filename="eq_chart.png"):
    try:
        # mood_history: [{"date": "2024-01-01", "score": 80}]
        df = pd.DataFrame(mood_history)
        plt.figure(figsize=(10, 6))
        sns.set_theme(style="whitegrid")
        plot = sns.lineplot(data=df, x="date", y="score", marker="o", color="#2DD4BF")
        plt.title("Zihinsel Verimlilik ve Duygu Analizi", fontsize=16, color="#0F172A")
        plt.ylim(0, 100)
        plt.xticks(rotation=45)
        
        output_path = os.path.join("/tmp", filename)
        plt.savefig(output_path, bbox_inches='tight')
        plt.close()
        return output_path
    except Exception as e:
        return f"Chart Error: {str(e)}"

def generate_ppt_code_snippet(slides, filename):
    """PowerPoint oluşturmak için kullanılan Python kodunu döndür"""
    num_slides = len(slides)
    
    code = f"""from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Sunucu oluşturma
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# {num_slides} slayt ekleme
slides_data = {repr([s.get('title', 'Başlık') for s in slides[:3]])}
...

for slide_data in slides_data:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Arka plan
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    # Başlık
    title = slide.shapes.title
    title.text = slide_data['title']
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(45, 212, 191)

# Kaydet
prs.save("/tmp/{filename}")
print(f"✅ {{len(prs.slides)}} slayt oluşturuldu")
"""
    return code

def generate_ppt_preview_image(slides, filename):
    """PowerPoint slaytlarının görsel önizlemesini oluştur"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from matplotlib.patches import Rectangle, FancyBboxPatch
        
        if not slides:
            return None
        
        # Gösterilecek slayt sayısı (max 3)
        n_slides = min(len(slides), 3)
        
        # Figür boyutları
        fig_width = 14
        fig_height = 3.5 * n_slides + 1
        
        fig, axes = plt.subplots(n_slides, 1, figsize=(fig_width, fig_height))
        if n_slides == 1:
            axes = [axes]
        
        for idx, (ax, slide_data) in enumerate(zip(axes, slides[:3])):
            ax.set_xlim(0, 10)
            ax.set_ylim(0, 5.625)  # 16:9 ratio
            ax.axis('off')
            
            # Slayt arka planı
            bg = FancyBboxPatch((0, 0), 10, 5.625, boxstyle="round,pad=0.02",
                                facecolor='#0F172A', edgecolor='#334155', linewidth=2)
            ax.add_patch(bg)
            
            # Başlık alanı
            title_bg = FancyBboxPatch((0.2, 4.5), 9.6, 0.8, boxstyle="round,pad=0.02",
                                       facecolor='#1E293B', edgecolor='#2DD4BF', linewidth=1, alpha=0.8)
            ax.add_patch(title_bg)
            
            # Başlık metni
            title_text = slide_data.get('title', 'Başlıksız')[:40]
            ax.text(5, 4.9, title_text, fontsize=11, fontweight='bold', 
                   color='#2DD4BF', ha='center', va='center')
            
            # İçerik alanı
            content = slide_data.get('content', [])
            y_pos = 3.8
            for point in content[:5]:  # Max 5 bullet
                # Bullet point
                circle = plt.Circle((0.5, y_pos), 0.08, color='#2DD4BF', alpha=0.6)
                ax.add_patch(circle)
                # Metin
                text = str(point)[:50] + ('...' if len(str(point)) > 50 else '')
                ax.text(0.8, y_pos, text, fontsize=8, color='#E2E8F0', 
                       va='center', ha='left')
                y_pos -= 0.6
            
            # Görsel placeholder (varsa)
            if slide_data.get('image_prompt') or slide_data.get('image_url'):
                img_rect = FancyBboxPatch((6.5, 1), 3, 3, boxstyle="round,pad=0.02",
                                          facecolor='#1E293B', edgecolor='#475569', 
                                          linewidth=1, linestyle='--')
                ax.add_patch(img_rect)
                ax.text(8, 2.5, '🖼️ Görsel', fontsize=9, color='#64748B', ha='center', va='center')
            
            # Slayt numarası
            ax.text(9.8, 0.2, f'{idx + 1}', fontsize=9, color='#64748B', ha='right', va='bottom')
        
        # Ana başlık
        fig.suptitle(f'📊 {filename} · {len(slides)} slayt', fontsize=14, 
                     fontweight='bold', color='#2DD4BF', y=0.98)
        
        # Görseli kaydet
        preview_path = os.path.join("/tmp", f"preview_{filename}.png")
        plt.savefig(preview_path, dpi=150, bbox_inches='tight', 
                   facecolor='#0F172A', edgecolor='none')
        plt.close()
        
        # Base64'e çevir
        with open(preview_path, 'rb') as f:
            img_data = base64.b64encode(f.read()).decode('utf-8')
        
        os.remove(preview_path)
        return img_data
        
    except Exception as e:
        print(f"PPT Preview image error: {e}", file=sys.stderr)
        return None

def generate_ppt(slides, filename):
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        for s_data in slides:
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(15, 23, 42)
            title = slide.shapes.title
            title.text = s_data.get("title", "Untitled")
            title_text = title.text_frame.paragraphs[0]
            title_text.font.bold = True
            title_text.font.size = Pt(36)
            title_text.font.color.rgb = RGBColor(45, 212, 191)
            tf = slide.placeholders[1].text_frame
            tf.text = "" 
            body = s_data.get("content", [])
            for point in body:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(226, 232, 240)
            prompt = s_data.get("image_prompt")
            image_url = s_data.get("image_url")
            actual_url = image_url
            if prompt:
                actual_url = f"https://image.pollinations.ai/prompt/{requests.utils.quote(prompt + ', high quality, concept art')}?width=1024&height=768&nologo=true"
            if actual_url:
                try:
                    response = requests.get(actual_url, timeout=10)
                    if response.status_code == 200:
                        img_path = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg").name
                        with open(img_path, 'wb') as f:
                            f.write(response.content)
                        left, top, width, height = Inches(7.5), Inches(1.5), Inches(5.3), Inches(5)
                        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                        os.unlink(img_path)
                except: pass
        output_path = os.path.join("/tmp", filename)
        prs.save(output_path)
        
        # Preview ve kod snippet oluştur
        code_snippet = generate_ppt_code_snippet(slides, filename)
        preview_image = generate_ppt_preview_image(slides, filename)
        
        return {
            "path": output_path,
            "code_snippet": code_snippet,
            "preview_image": preview_image,
            "slides": len(slides)
        }
    except Exception as e:
        return f"Error: {str(e)}"

def main():
    try:
        input_data = sys.stdin.read()
        if not input_data: return
        payload = json.loads(input_data)
        mode = payload.get("mode") # "export" or "search"
        
        if mode == "search":
            query = payload.get("query")
            print(json.dumps({"results": web_search(query)}))
            return

        file_type = payload.get("type")
        filename = payload.get("filename", "exported_file")
        
        if file_type == "excel":
            result = generate_excel(payload.get("data", {}), filename)
            # Yeni format: result bir dict (path, code_snippet, preview_image, rows, columns)
            if isinstance(result, dict):
                print(json.dumps({
                    "success": True,
                    "path": result["path"],
                    "code_snippet": result.get("code_snippet"),
                    "preview_image": result.get("preview_image"),
                    "rows": result.get("rows"),
                    "columns": result.get("columns")
                }))
            elif isinstance(result, str) and result.startswith("Error:"):
                print(json.dumps({"error": result}))
            else:
                print(json.dumps({"success": True, "path": result}))
        elif file_type == "word":
            content = payload.get("content", "")
            # If EQ Analysis, generate chart first
            if payload.get("eq_data"):
                chart_path = generate_eq_chart(payload["eq_data"], "temp_eq.png")
                # Insert chart into list of content
                content = [{"type": "heading", "text": "Duygusal Zeka ve Analiz Raporu", "level": 0}] + \
                          content + \
                          [{"type": "heading", "text": "Haftalık İlerleme Grafiği", "level": 1},
                           {"type": "image", "path": chart_path}]
            result = generate_word(content, filename)
            if isinstance(result, str) and result.startswith("Error:"):
                print(json.dumps({"error": result}))
            else:
                print(json.dumps({"success": True, "path": result}))
        elif file_type == "ppt":
            result = generate_ppt(payload.get("slides", []), filename)
            # Yeni format: result bir dict (path, code_snippet, preview_image, slides)
            if isinstance(result, dict):
                print(json.dumps({
                    "success": True,
                    "path": result["path"],
                    "code_snippet": result.get("code_snippet"),
                    "preview_image": result.get("preview_image"),
                    "slides": result.get("slides")
                }))
            elif isinstance(result, str) and result.startswith("Error:"):
                print(json.dumps({"error": result}))
            else:
                print(json.dumps({"success": True, "path": result}))
        elif file_type == "certificate":
            result = generate_certificate(payload.get("data", {}), filename)
            if isinstance(result, str) and result.startswith("Error:"):
                print(json.dumps({"error": result}))
            else:
                print(json.dumps({"success": True, "path": result}))
        else: 
            print(json.dumps({"error": "Error: Invalid file type"}))
        
    except Exception as e: 
        print(json.dumps({"error": str(e)}))

if __name__ == "__main__":
    main()
